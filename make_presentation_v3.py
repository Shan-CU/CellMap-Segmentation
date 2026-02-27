#!/usr/bin/env python3
"""
make_presentation_v3.py — CellMap Lab Meeting Presentation (Phase 2 Focus)

Revised narrative — no BST details, focused on Phase 2 architecture comparison.

Slides:
  1. Title
  2. The Big Picture: Why CellMap
  3. The Challenge: Partial Annotation & Class Imbalance
  4. The Roadmap (8-step pipeline, updated progress)
  5. Phase 1: What 63 Ablation Experiments Taught Us
  6. The Optimal Recipe — How We Got Here
  7. Phase 2: The Architecture Zoo (4×2D + 5×3D)
  8. Phase 2: Current Status & Early Results
  9. What's Next — Threshold Tuning → Ensemble → Novel Tomograms
 10. Thank You / Questions

Uses UNC lab template: GStephenson_Rotation3_Slides.pptx
"""

import io
import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# ── Paths ──────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE   = os.path.join(SCRIPT_DIR, "GStephenson_Rotation3_Slides.pptx")
FIG_DIR    = os.path.join(SCRIPT_DIR, "figures")

# ── Colour palette (matches template) ─────────────────────────────────
GOLD    = RGBColor(0xCF, 0xB8, 0x7C)
BLACK   = RGBColor(0x00, 0x00, 0x00)
DKGRAY  = RGBColor(0x33, 0x33, 0x33)
GRAY    = RGBColor(0x76, 0x76, 0x76)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x15, 0x60, 0x82)
ORANGE  = RGBColor(0xE9, 0x71, 0x32)
GREEN   = RGBColor(0x19, 0x6B, 0x24)
RED     = RGBColor(0xC0, 0x39, 0x2B)
LTGRAY  = RGBColor(0xD9, 0xD9, 0xD9)
LTBLUE  = RGBColor(0xD0, 0xDE, 0xE8)
CREAM   = RGBColor(0xF5, 0xF0, 0xE6)
LTGREEN = RGBColor(0xD4, 0xE8, 0xD4)
ZEBRA   = RGBColor(0xF2, 0xF2, 0xF2)
TAN     = RGBColor(0xE8, 0xE0, 0xD0)

# ── Helper functions ───────────────────────────────────────────────────

def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_figure(slide, name, left, top, max_width, max_height):
    from PIL import Image
    path = os.path.join(FIG_DIR, name)
    if not os.path.exists(path):
        print(f"  \u26a0 Figure not found: {path}")
        return
    with Image.open(path) as img:
        w_px, h_px = img.size
    aspect = w_px / h_px
    w = max_width
    h = int(w / aspect)
    if h > max_height:
        h = max_height
        w = int(h * aspect)
    slide.shapes.add_picture(path, left, top, w, h)


def set_para(tf, text, size=18, bold=False, color=DKGRAY, alignment=PP_ALIGN.LEFT,
             space_after=Pt(4), space_before=Pt(0), font_name="Calibri"):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def add_para(tf, text, size=18, bold=False, color=DKGRAY, alignment=PP_ALIGN.LEFT,
             space_after=Pt(4), space_before=Pt(0), font_name="Calibri", level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    p.level = level
    return p


def add_bullet_list(tf, items, size=16, color=DKGRAY, bold=False,
                    space_after=Pt(6), level=0, bullet="\u2022"):
    for item in items:
        add_para(tf, f"{bullet} {item}", size=size, color=color, bold=bold,
                 space_after=space_after, level=level)


def add_table(slide, left, top, width, height, rows,
              header_bg=BLUE, header_fg=WHITE, font_size=11):
    n_rows = len(rows)
    n_cols = len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.font.name = "Calibri"
                if r_idx == 0:
                    para.font.bold = True
                    para.font.color.rgb = header_fg
                else:
                    para.font.color.rgb = DKGRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ZEBRA
    return tbl


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ── Load template, extract logo, delete template slides ────────────────
_tpl = Presentation(TEMPLATE)

logo_blob = None
logo_ext  = "png"
for shape in _tpl.slides[1].shapes:
    if shape.name == "Picture 1":
        logo_blob = shape.image.blob
        logo_ext  = shape.image.content_type.split("/")[-1]
        break


def _delete_slide(prs, slide_index):
    rIds = []
    slide_part = prs.slides[slide_index].part
    for rId, rel in prs.part.rels.items():
        if rel.target_part is slide_part:
            rIds.append(rId)
    sldIdLst = prs.slides._sldIdLst
    ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    for sldId in list(sldIdLst):
        if sldId.get(ns + 'id') in rIds:
            sldIdLst.remove(sldId)
    for rId in rIds:
        prs.part.rels._rels.pop(rId, None)


for i in range(len(_tpl.slides) - 1, -1, -1):
    _delete_slide(_tpl, i)

prs = _tpl
title_layout = prs.slide_layouts[0]


def add_content_slide(title_text):
    slide = prs.slides.add_slide(title_layout)
    # White background
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0),
                                   Inches(13.33), Inches(6.50))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.fill.background()
    # Title
    tf = add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.55))
    set_para(tf, title_text, size=28, bold=True, color=GOLD, font_name="Calibri Light")
    # Gold rule
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), Inches(0.67),
                                   Inches(12.33), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    # Logo
    if logo_blob:
        buf = io.BytesIO(logo_blob)
        slide.shapes.add_picture(buf, Inches(11.45), Inches(6.57),
                                 Inches(1.88), Inches(0.86))
    return slide


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(title_layout)

tf = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10.93), Inches(2.0))
set_para(tf, "Multi-Label Organelle Segmentation", size=42, bold=True,
         color=GOLD, alignment=PP_ALIGN.LEFT, font_name="Calibri Light")
add_para(tf, "from FIB-SEM Electron Microscopy", size=42, bold=True,
         color=GOLD, alignment=PP_ALIGN.LEFT, font_name="Calibri Light")

tf = add_textbox(slide, Inches(1.2), Inches(3.8), Inches(10.93), Inches(0.6))
set_para(tf, "Toward Automated Annotation of Novel Electron Tomograms",
         size=22, color=DKGRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri")

tf = add_textbox(slide, Inches(1.2), Inches(5.0), Inches(10.93), Inches(0.5))
set_para(tf, "George S. George", size=24, bold=False, color=DKGRAY,
         alignment=PP_ALIGN.LEFT, font_name="Calibri")

tf = add_textbox(slide, Inches(1.2), Inches(5.5), Inches(10.93), Inches(0.5))
set_para(tf, "Department of Computer Science  \u2022  Interdisciplinary Quantitative Biology Program",
         size=14, color=GRAY, alignment=PP_ALIGN.LEFT)

tf = add_textbox(slide, Inches(1.2), Inches(5.85), Inches(10.93), Inches(0.3))
set_para(tf, "Lab Meeting  \u2022  February 2026", size=14, color=GRAY,
         alignment=PP_ALIGN.LEFT)

set_notes(slide,
    "Hi everyone, thanks for having me. "
    "Today I'll walk you through my rotation project \u2014 building a segmentation "
    "model for electron microscopy data. The long-term goal is a model that can take "
    "novel electron tomograms and automatically label every organelle.\n\n"
    "I'll cover: the CellMap challenge and why it's our proving ground, "
    "the key challenges \u2014 partial annotation and extreme class imbalance, "
    "what we learned from 63 ablation experiments, the optimal training recipe, "
    "and the Phase 2 architecture comparison that's running right now."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE BIG PICTURE — WHY CELLMAP
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("The Big Picture: Why CellMap?")

# Left column — goal and challenges
tf = add_textbox(slide, Inches(0.5), Inches(0.85), Inches(6.0), Inches(5.3))
set_para(tf, "The Long-Term Goal", size=20, bold=True, color=BLUE)
add_bullet_list(tf, [
    "Build a model that takes novel electron\n"
    "tomograms and produces fully labeled\n"
    "organelle segmentations \u2014 automatically",
    "No manual annotation \u2014 the model\n"
    "generalizes to unseen biology",
    "Everything we learn here transfers\n"
    "directly to our own tomograms",
], size=14, color=DKGRAY)

add_para(tf, "", size=10)
add_para(tf, "Why Is This Hard?", size=20, bold=True, color=BLUE, space_before=Pt(8))
add_bullet_list(tf, [
    "48 organelle classes, 22 biological samples",
    "Partial annotations \u2014 unlabeled \u2260 absent;\n"
    "standard losses penalize correct predictions",
    "6+ orders of magnitude class imbalance",
    "3D volumes: up to 512 GB host RAM just\n"
    "for data loading",
], size=14, color=DKGRAY)

# Right: CellMap proving ground box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7.0), Inches(0.85), Inches(5.8), Inches(2.5))
rect.fill.solid()
rect.fill.fore_color.rgb = CREAM
rect.line.fill.background()

tf2 = add_textbox(slide, Inches(7.2), Inches(0.95), Inches(5.4), Inches(2.3))
set_para(tf2, "Why CellMap as the Training Ground?", size=16, bold=True, color=GOLD)
add_bullet_list(tf2, [
    "Largest public FIB-SEM segmentation benchmark",
    "Diverse biology: HeLa, mouse kidney/liver/\n"
    "heart, fly brain, jurkat, macrophage\u2026",
    "Has a leaderboard \u2192 objective evaluation",
    "Same imaging modality as our own data",
], size=13, color=DKGRAY)

# Key stats boxes
data = [
    ("48", "organelle classes"),
    ("22", "biological samples"),
    ("289", "annotated 3D crops"),
    ("42.3B", "total voxels"),
]
for i, (num, label) in enumerate(data):
    y = Inches(3.7) + Inches(i * 0.75)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.2), y, Inches(5.4), Inches(0.65))
    rect.fill.solid()
    rect.fill.fore_color.rgb = TAN
    rect.line.fill.background()
    tf3 = add_textbox(slide, Inches(7.4), y + Inches(0.05), Inches(1.5), Inches(0.55))
    set_para(tf3, num, size=26, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    tf4 = add_textbox(slide, Inches(8.8), y + Inches(0.12), Inches(3.6), Inches(0.4))
    set_para(tf4, label, size=14, color=GRAY)

set_notes(slide,
    "The long-term goal is to build a model that takes in novel electron tomograms "
    "and automatically produces fully labeled organelle segmentations.\n\n"
    "Why CellMap? It's the largest publicly available FIB-SEM segmentation benchmark. "
    "48 organelle classes, 22 biological samples spanning HeLa cells, mouse organs, "
    "fly brain, and more. It has a leaderboard so we can objectively measure progress. "
    "And critically, it's the same imaging modality as our own tomograms, so everything "
    "we learn transfers directly."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE CHALLENGE — PARTIAL ANNOTATION & IMBALANCE
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("The Challenge: Partial Annotation & Class Imbalance")

# Left: Partial annotation explanation
tf = add_textbox(slide, Inches(0.5), Inches(0.85), Inches(6.0), Inches(2.5))
set_para(tf, "Partial Annotation", size=20, bold=True, color=BLUE)
add_para(tf, "", size=6)
add_bullet_list(tf, [
    "Each crop only labels a subset of 48 classes",
    'Unlabeled voxels are NaN \u2014 NOT "absent"',
    "Na\u00efve loss treats NaN as negative \u2192\n"
    "penalizes correct predictions",
    "Solution: NaN-masked loss \u2014 only compute\n"
    "gradients where annotations exist",
], size=14, color=DKGRAY)

# Right: Class imbalance explanation
tf = add_textbox(slide, Inches(6.8), Inches(0.85), Inches(6.0), Inches(2.5))
set_para(tf, "Extreme Class Imbalance", size=20, bold=True, color=BLUE)
add_para(tf, "", size=6)
add_bullet_list(tf, [
    "6+ orders of magnitude between largest\n"
    "and smallest classes",
    '"cell" = billions of voxels;\n'
    '"mito_ribo" = tens of thousands',
    "Without mitigation, model only learns\n"
    "the top 5\u201310 largest classes",
    "Solution: weighted sampling + foreground\n"
    "masking to focus on annotated regions",
], size=14, color=DKGRAY)

# Bottom: Visual example — class distribution
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(3.7), Inches(12.3), Inches(2.7))
rect.fill.solid()
rect.fill.fore_color.rgb = ZEBRA
rect.line.fill.background()

tf = add_textbox(slide, Inches(0.7), Inches(3.80), Inches(11.9), Inches(0.4))
set_para(tf, "Class Hierarchy: 48 Classes = 31 Atomic + 17 Group/Composite",
         size=16, bold=True, color=BLUE)

# Organelle categories
categories = [
    ("Mitochondria", "mito_mem, mito_lum, mito_ribo \u2192 mito", BLUE),
    ("ER", "er_mem, er_lum, eres_mem, eres_lum \u2192 er, eres, er_mem_all", BLUE),
    ("Nucleus", "ne_mem, ne_lum, np_out, np_in, hchrom, echrom, nucpl \u2192 nuc, ne, np, chrom", BLUE),
    ("Vesicles", "ves_mem, ves_lum, endo_mem, endo_lum, lyso_mem, lyso_lum \u2192 ves, endo, lyso", BLUE),
    ("Other", "golgi_mem/lum, ld_mem/lum, perox_mem/lum, mt_out/in, ecs, pm, cyto \u2192 cell", GRAY),
]

for i, (cat, members, clr) in enumerate(categories):
    y = Inches(4.25) + Inches(i * 0.42)
    tf2 = add_textbox(slide, Inches(0.9), y, Inches(2.0), Inches(0.38))
    set_para(tf2, cat, size=12, bold=True, color=clr)
    tf3 = add_textbox(slide, Inches(2.9), y, Inches(9.5), Inches(0.38))
    set_para(tf3, members, size=10, color=GRAY, font_name="Consolas")

set_notes(slide,
    "Two core challenges dominate this problem.\n\n"
    "First, partial annotation. Each 3D crop only labels a subset of the 48 classes. "
    "Unlabeled voxels are NaN, meaning 'not annotated' \u2014 NOT 'absent.' If you "
    "na\u00efvely treat NaN as negative, you penalize the model for correct predictions. "
    "We solve this with NaN-masked loss: gradients only flow where annotations exist.\n\n"
    "Second, extreme class imbalance \u2014 over 6 orders of magnitude. The 'cell' class "
    "has billions of voxels; 'mito_ribo' has tens of thousands. Without mitigation, "
    "the model only learns the biggest 5\u201310 classes and ignores everything else. "
    "We use weighted sampling and foreground masking to address this.\n\n"
    "The 48 classes form a hierarchy: 31 atomic labels like 'mito_mem' for mitochondrial "
    "membrane, plus 17 group labels like 'mito' = union of mito_mem + mito_lum + mito_ribo."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: THE ROADMAP (updated progress)
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("The Roadmap")

steps = [
    ("1", "Infrastructure",          "Native zarr pipeline, all 48 classes, partial annotation support",  "done"),
    ("2", "Ablation Study",          "63 experiments across 5 sweeps (loss, \u03b1\u03b2, \u03c4, masking, techniques)", "done"),
    ("3", "Compose Optimal Config",  "Winner from each sweep \u2192 single recipe",                      "done"),
    ("4", "Architecture Comparison", "4\u00d72D + 5\u00d73D models, full training with Dice tracking",   "here"),
    ("5", "Threshold Tuning",        "Per-class sigmoid optimization on validation",                      ""),
    ("6", "Per-Class Ensemble",      "Best architecture per organelle \u2192 composite model",            ""),
    ("7", "Leaderboard Submission",  "Submit to CellMap challenge",                                       ""),
    ("8", "Novel Tomograms",         "Apply to our own electron tomography data",                         "goal"),
]

y_start = Inches(0.90)
row_h   = Inches(0.62)
gap     = Inches(0.06)

for i, (num, title, desc, status) in enumerate(steps):
    y = y_start + (row_h + gap) * i

    if status == "done":
        bg = TAN
        num_bg = GOLD
    elif status == "here":
        bg = LTBLUE
        num_bg = BLUE
    elif status == "goal":
        bg = LTGREEN
        num_bg = GOLD
    else:
        bg = ZEBRA
        num_bg = LTGRAY

    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(0.65), y + Inches(0.06),
                                   Inches(0.48), Inches(0.48))
    circ.fill.solid()
    circ.fill.fore_color.rgb = num_bg
    circ.line.fill.background()
    tf_n = add_textbox(slide, Inches(0.65), y + Inches(0.09),
                        Inches(0.48), Inches(0.44))
    set_para(tf_n, num, size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Step box
    box_w = Inches(3.8)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.30), y, box_w, row_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    if status == "here":
        rect.line.color.rgb = BLUE
        rect.line.width = Pt(2)

    tf_t = add_textbox(slide, Inches(1.50), y + Inches(0.08), Inches(3.4), Inches(0.44))
    title_color = BLUE if status == "here" else DKGRAY
    set_para(tf_t, title, size=15, bold=True, color=title_color)

    tf_d = add_textbox(slide, Inches(5.30), y + Inches(0.08), Inches(5.0), Inches(0.44))
    set_para(tf_d, desc, size=13, color=GRAY)

    if status == "done":
        tf_s = add_textbox(slide, Inches(10.5), y + Inches(0.08), Inches(2.3), Inches(0.44))
        set_para(tf_s, "\u2705 Done", size=14, bold=True, color=GREEN,
                 alignment=PP_ALIGN.RIGHT)
    elif status == "here":
        tf_s = add_textbox(slide, Inches(10.5), y + Inches(0.08), Inches(2.3), Inches(0.44))
        set_para(tf_s, "\u2190 In Progress", size=14, bold=True, color=BLUE,
                 alignment=PP_ALIGN.RIGHT)

    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                        Inches(3.0), y + row_h,
                                        Inches(0.30), gap)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

set_notes(slide,
    "Here's the full pipeline. Steps 1 through 3 are complete.\n\n"
    "Step 1: We built the infrastructure \u2014 native zarr data loading, all 48 classes, "
    "partial annotation masking.\n"
    "Step 2: 63 ablation experiments across 5 sweeps to find the optimal training recipe.\n"
    "Step 3: Composed the winning configuration from each sweep.\n"
    "Step 4 is where we are now \u2014 training 9 architectures (4 in 2D, 5 in 3D) with "
    "the optimal recipe. This is running on Longleaf L40S GPUs right now.\n\n"
    "Steps 5\u20138 are ahead: per-class threshold tuning, per-class ensemble, "
    "leaderboard submission, and finally applying to our own novel tomograms."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5: PHASE 1 — WHAT 63 ABLATION EXPERIMENTS TAUGHT US
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("Phase 1: What 63 Ablation Experiments Taught Us")

# Sweep summary at top
tf = add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4))
set_para(tf, "34 in 2D (L40S, ~3h each)  +  29 in 3D (H100, ~15h each)  =  63 experiments, ~600 GPU-hours",
         size=13, color=GRAY)

# Results table — compact, focused on findings not individual experiments
sweep_rows = [
    ["Sweep", "What We Varied", "# Exps", "Key Finding", "Winner"],
    ["A: Loss", "8 loss functions\n(BCE, Dice, Tversky, Focal\u2026)", "16",
     "Simple losses crush complex ones.\nBCE: 0.476 Dice, DiceBCE: 0.459 Dice",
     "dice_bce"],
    ["B: Tversky \u03b1\u03b2", "6 precision/recall\nbias settings", "12",
     "All Tversky variants cluster at 0.37\u20130.41.\nNo \u03b1\u03b2 beats dice_bce.",
     "N/A \u2014 Tversky\nsuboptimal"],
    ["C: Class wt \u03c4", "5 logit adjustment\nstrengths", "10",
     "High \u03c4 destroys predictions (Dice\u21920).\nLogit adjustment + partial annot = bad.",
     "N/A \u2014 no logit\nadjustment"],
    ["D: Masking", "7 masking strategies\n(FG, bbox, masksup\u2026)", "14",
     "FG mask gives free +28% improvement.\nBbox alone HURTS performance.",
     "FG mask ON,\nbbox OFF"],
    ["E: Techniques", "EMA, sampler, aug,\ncrop weights", "11",
     "EMA: 4\u00d7 better val_loss. Sampler: essential.\nIntensity aug HURTS (\u221223%). Crop wts HURT (\u221238%).",
     "EMA + weighted\nsampler"],
]
add_table(slide, Inches(0.3), Inches(1.30), Inches(12.7), Inches(3.8),
          sweep_rows, font_size=10, header_bg=BLUE)

# Bottom: key takeaways
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.1))
rect.fill.solid()
rect.fill.fore_color.rgb = LTBLUE
rect.line.fill.background()

tf = add_textbox(slide, Inches(0.7), Inches(5.30), Inches(11.9), Inches(1.0))
set_para(tf, "Key Takeaway:", size=14, bold=True, color=BLUE)
add_para(tf, "Simple beats complex. DiceBCE + EMA + foreground mask + weighted sampler = best recipe.  "
         "Fancy losses (Tversky, focal, logit adjustment) all underperformed. "
         "Extra augmentation and sampling tricks actively hurt.",
         size=13, color=DKGRAY)

set_notes(slide,
    "63 ablation experiments across 5 sweeps. Here are the key findings.\n\n"
    "Sweep A, loss function: we tested 8 losses. Simple losses like BCE and DiceBCE "
    "massively outperformed complex ones like Tversky and focal variants. The best "
    "leaderboard Dice was BCE at 0.476, but we chose DiceBCE at 0.459 because it has "
    "better rare-class performance \u2014 the Dice component explicitly optimizes overlap.\n\n"
    "Sweep B and C: Tversky alpha-beta tuning and logit adjustment both failed to "
    "improve on DiceBCE. In fact, high logit adjustment completely destroyed predictions.\n\n"
    "Sweep D: foreground masking gave a free 28% improvement. But bounding box masking "
    "alone actually hurt \u2014 surprising result.\n\n"
    "Sweep E: EMA gave a 4x improvement in val_loss. Weighted sampler is essential. "
    "But intensity augmentation hurt by 23% and class-aware crop weighting hurt by 38%. "
    "EM data has very consistent intensity within volumes, so augmenting it adds noise "
    "rather than useful variation.\n\n"
    "The bottom line: simple beats complex. Our optimal recipe is DiceBCE + EMA + "
    "foreground mask + weighted sampler. That's it."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6: THE OPTIMAL RECIPE — HOW WE GOT HERE
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("The Optimal Recipe")

# Central config box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(0.85), Inches(5.8), Inches(5.3))
rect.fill.solid()
rect.fill.fore_color.rgb = CREAM
rect.line.fill.background()

tf = add_textbox(slide, Inches(0.7), Inches(0.95), Inches(5.4), Inches(5.1))
set_para(tf, "Phase 2 Training Configuration", size=18, bold=True, color=GOLD)
add_para(tf, "", size=8)

config_items = [
    ("Loss", "DiceBCE", "Dice component optimizes\noverlap for rare classes"),
    ("EMA", "Enabled (decay=0.999)", "Smooths noisy gradients from\npartial 48-class annotations"),
    ("FG Mask", "Enabled", "Focus loss on foreground voxels;\nfree +28% improvement"),
    ("Weighted Sampler", "Enabled", "cellmap-data\u2019s default;\nessential for class balance"),
    ("Intensity Aug", "DISABLED", "Hurt \u221223% \u2014 EM intensity\nis already consistent"),
    ("Crop Weighting", "DISABLED", "Hurt \u221238% \u2014 disrupts\nbalanced sampling"),
    ("Deep Supervision", "SegResNet only", "Multi-scale auxiliary losses\nat decoder layers"),
    ("AMP", "Enabled", "Mixed precision for speed"),
    ("Scheduler", "Cosine + warmup", "10-epoch warmup, cosine decay"),
    ("Optimizer", "RAdam, lr=1e-4", "Robust to LR sensitivity"),
]

for param, value, reason in config_items:
    add_para(tf, f"{param}:  {value}", size=12, bold=True, color=BLUE, space_after=Pt(1))
    add_para(tf, f"     {reason}", size=10, color=GRAY, space_after=Pt(8))

# Right side: Why this works / 3D note
tf = add_textbox(slide, Inches(6.6), Inches(0.85), Inches(6.2), Inches(2.8))
set_para(tf, "Why DiceBCE Over BCE?", size=18, bold=True, color=BLUE)
add_para(tf, "", size=6)
add_bullet_list(tf, [
    "BCE had highest mean Dice (0.476) but\n"
    "achieves it by over-predicting common classes",
    "DiceBCE (0.459) has better rare-class\n"
    "performance \u2014 Dice component explicitly\n"
    "optimizes per-class overlap",
    "Challenge weights all 48 classes equally,\n"
    "so rare-class performance matters",
    "With EMA, DiceBCE reaches val_loss = 0.112\n"
    "\u2014 the best result by far (4\u00d7 better\n"
    "than without EMA)",
], size=13, color=DKGRAY)

# 3D note
rect2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.6), Inches(3.9), Inches(6.2), Inches(2.2))
rect2.fill.solid()
rect2.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xE0)
rect2.line.fill.background()

tf = add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.0))
set_para(tf, "\u26a0  3D Ablation Was Non-Functional", size=15, bold=True, color=ORANGE)
add_para(tf, "", size=6)
add_bullet_list(tf, [
    "All 29 3D experiments: max Dice = 0.017",
    "12,500 training steps was woefully\n"
    "insufficient for 128\u00b3 volumes",
    "Only 'nuc' (largest organelle) got any signal",
    "Decision: transfer 2D findings to 3D,\n"
    "train 10\u201324\u00d7 longer in Phase 2\n"
    "(300K steps vs 12.5K)",
], size=12, color=DKGRAY)

set_notes(slide,
    "Here's the exact configuration going into Phase 2.\n\n"
    "Why DiceBCE over BCE? BCE had the highest raw Dice at 0.476, but it achieves that "
    "by over-predicting common classes like ecs, pm, and mito. DiceBCE at 0.459 has "
    "better rare-class performance because the Dice component explicitly optimizes "
    "per-class overlap. Since the challenge weights all 48 classes equally, rare-class "
    "performance matters more than it looks.\n\n"
    "With EMA, DiceBCE reaches a val_loss of 0.112 \u2014 4x better than without. "
    "EMA smooths noisy gradients from the partial annotation setup where different "
    "crops annotate different subsets of classes.\n\n"
    "Important note: our 3D ablation experiments were essentially non-functional. "
    "Every experiment produced near-zero Dice. The training regime was only 12,500 steps "
    "\u2014 way too few for 3D 128-cubed volumes. So we're transferring the 2D findings "
    "to 3D and training 24 times longer in Phase 2: 300,000 steps."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7: PHASE 2 — THE ARCHITECTURE ZOO
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("Phase 2: The Architecture Zoo")

tf = add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.3))
set_para(tf, "Same recipe, 9 architectures \u2014 which backbone learns organelles best?",
         size=16, color=GRAY)

# 2D models table
tf_2d = add_textbox(slide, Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.35))
set_para(tf_2d, "2D Models (slice-by-slice, 256\u00d7256)", size=16, bold=True, color=BLUE)

rows_2d = [
    ["Architecture", "Params", "Type", "Strength"],
    ["ResNet Generator", "7.8M", "CNN", "Fast, lightweight encoder-decoder baseline"],
    ["UNet", "31M", "CNN", "Classic encoder-decoder, large capacity"],
    ["SwinTransformer", "36M", "Transformer", "Shifted windows, multi-scale attention"],
    ["ViTVNet", "105M", "Hybrid", "ViT encoder + V-Net decoder (largest)"],
]
add_table(slide, Inches(0.3), Inches(1.65), Inches(6.3), Inches(2.0),
          rows_2d, font_size=11, header_bg=BLUE)

# 3D models table
tf_3d = add_textbox(slide, Inches(0.5), Inches(3.85), Inches(6.0), Inches(0.35))
set_para(tf_3d, "3D Models (full volumes, 128\u00b3)", size=16, bold=True, color=BLUE)

rows_3d = [
    ["Architecture", "Params", "Type", "Strength"],
    ["SegResNetDS", "19.9M", "CNN", "Deep supervision, multi-scale loss"],
    ["SwinUNETR", "62.2M", "Transformer", "3D shifted-window self-attention"],
    ["ResNet3D", "24.5M", "CNN", "3D version of CSC ResNet generator"],
    ["UNet3D", "90.3M", "CNN", "Largest 3D model, pure encoder-decoder"],
    ["ViTVNet3D", "31.5M", "Hybrid", "ViT encoder + V-Net decoder in 3D"],
]
add_table(slide, Inches(0.3), Inches(4.25), Inches(6.3), Inches(2.2),
          rows_3d, font_size=11, header_bg=BLUE)

# Right side: training config comparison
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7.0), Inches(1.25), Inches(5.8), Inches(2.4))
rect.fill.solid()
rect.fill.fore_color.rgb = CREAM
rect.line.fill.background()

tf = add_textbox(slide, Inches(7.2), Inches(1.35), Inches(5.4), Inches(2.2))
set_para(tf, "2D Training", size=14, bold=True, color=GOLD)
add_para(tf, "100 epochs \u00d7 1,000 iter = 100K steps", size=12, color=DKGRAY)
add_para(tf, "Batch size 8, single L40S GPU", size=12, color=DKGRAY)
add_para(tf, "~6\u201310 hours per model", size=12, color=GRAY)
add_para(tf, "Validate every 10 epochs with Dice", size=12, color=DKGRAY)

add_para(tf, "", size=8)
add_para(tf, "3D Training", size=14, bold=True, color=GOLD)
add_para(tf, "1,000 epochs \u00d7 300 iter = 300K steps", size=12, color=DKGRAY)
add_para(tf, "Batch size 2, single L40S (48GB VRAM)", size=12, color=DKGRAY)
add_para(tf, "~10\u201313 days per model", size=12, color=GRAY)
add_para(tf, "Validate every 30 epochs with Dice", size=12, color=DKGRAY)
add_para(tf, "persistent_workers=false (OOM fix)", size=12, color=DKGRAY)

# Why 2D AND 3D?
rect2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(7.0), Inches(3.85), Inches(5.8), Inches(2.6))
rect2.fill.solid()
rect2.fill.fore_color.rgb = LTBLUE
rect2.line.fill.background()

tf = add_textbox(slide, Inches(7.2), Inches(3.95), Inches(5.4), Inches(2.4))
set_para(tf, "Why Both 2D and 3D?", size=14, bold=True, color=BLUE)
add_para(tf, "", size=6)
add_bullet_list(tf, [
    "2D: fast iteration, strong on in-plane\n"
    "features (membranes, texture)",
    "3D: captures volumetric context\n"
    "(tubules, 3D connectivity, shape)",
    "Per-class ensemble: some organelles\n"
    "may be better served by 2D, others\n"
    "by 3D \u2192 composite model",
    "Final ensemble selects best architecture\n"
    "per organelle class",
], size=12, color=DKGRAY)

set_notes(slide,
    "Phase 2: 9 architectures, same recipe. The question is which backbone learns "
    "organelles best.\n\n"
    "4 in 2D: ResNet (our Phase 1 workhorse at 7.8M params), UNet at 31M, "
    "SwinTransformer at 36M with shifted-window attention, and ViTVNet at 105M \u2014 "
    "our largest model, using a Vision Transformer encoder.\n\n"
    "5 in 3D: SegResNetDS at 20M with deep supervision, SwinUNETR at 62M with 3D "
    "self-attention, ResNet3D at 24.5M, UNet3D at 90M (largest 3D model), and "
    "ViTVNet3D at 31.5M.\n\n"
    "2D models process 256x256 slices independently \u2014 fast but no inter-slice context. "
    "3D models process 128-cubed volumes \u2014 capture tubular structures and 3D connectivity "
    "but are 30x more expensive per step. All models run on single L40S GPUs with 48GB VRAM.\n\n"
    "The final ensemble will select the best architecture per organelle class. Some "
    "organelles like flat membranes may be better served by 2D; others like tubular ER "
    "or microtubules may need 3D context."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: PHASE 2 — CURRENT STATUS & EARLY RESULTS
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("Phase 2: Current Status")

# Status table
status_rows = [
    ["Model", "GPU", "Batch", "Steps", "Val Every", "Status"],
    ["ResNet 2D", "L40S", "8", "100K", "10 ep", "\u23f3 Pending"],
    ["UNet 2D", "L40S", "8", "100K", "10 ep", "\u23f3 Pending"],
    ["Swin 2D", "L40S", "8", "100K", "10 ep", "\u23f3 Pending"],
    ["ViT 2D", "L40S", "8", "100K", "10 ep", "\u23f3 Pending"],
    ["SegResNet 3D", "L40S", "2", "300K", "30 ep", "\u23f3 Pending"],
    ["SwinUNETR 3D", "L40S", "2", "300K", "30 ep", "\u23f3 Pending"],
    ["ResNet 3D", "L40S", "2", "300K", "30 ep", "\u23f3 Pending"],
    ["UNet 3D", "L40S", "2", "300K", "30 ep", "\u23f3 Pending"],
    ["ViTNet 3D", "L40S", "1", "300K", "30 ep", "\u23f3 Pending"],
]
add_table(slide, Inches(0.3), Inches(0.85), Inches(8.0), Inches(3.5),
          status_rows, font_size=10, header_bg=BLUE)

# Right side: what's different this time
tf = add_textbox(slide, Inches(8.5), Inches(0.85), Inches(4.5), Inches(3.5))
set_para(tf, "What\u2019s Different This Time", size=16, bold=True, color=BLUE)
add_para(tf, "", size=6)
add_bullet_list(tf, [
    "Clean start: all old checkpoints\n"
    "and outputs wiped",
    "Per-class Dice from epoch 1:\n"
    "48 classes tracked in TensorBoard",
    "Best model saved by val_dice\n"
    "(not val_loss)",
    "Fixed checkpoint resume bug:\n"
    "best metrics now survive restarts",
    "3D OOM fix: persistent_workers=\n"
    "false prevents memory leak",
], size=12, color=DKGRAY)

# Bottom: infrastructure challenges
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.3), Inches(4.55), Inches(12.7), Inches(1.9))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xE0)
rect.line.fill.background()

tf = add_textbox(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(1.7))
set_para(tf, "Infrastructure Challenges Overcome", size=14, bold=True, color=ORANGE)
add_para(tf, "", size=4)

# Four columns of challenges
challenges = [
    ("3D Host RAM OOM", "CellMapDataLoader.refresh() leaks\nmemory via persistent workers.\nFix: persistent_workers=false"),
    ("No Dice Tracking", "Original code only logged val_loss.\nAdded per-class Dice for all 48\nclasses with NaN masking."),
    ("Checkpoint Bug", "best_val_dice was reset after\ncheckpoint restore, losing best\nmodel on job preemption. Fixed."),
    ("MONAI API Break", "SwinUNETR removed img_size in\nv1.5.2. ViT patch_size fix.\nBoth resolved."),
]

for i, (title, desc) in enumerate(challenges):
    x = Inches(0.3) + Inches(i * 3.2)
    tf2 = add_textbox(slide, x, Inches(4.95), Inches(3.0), Inches(0.3))
    set_para(tf2, title, size=11, bold=True, color=ORANGE)
    tf3 = add_textbox(slide, x, Inches(5.2), Inches(3.0), Inches(1.1))
    set_para(tf3, desc, size=10, color=GRAY)

set_notes(slide,
    "We just launched a completely fresh Phase 2 run \u2014 all 9 models plus TensorBoard, "
    "all from scratch with clean outputs and optimized configurations.\n\n"
    "Every model gets an identical, fair setup: DiceBCE loss, EMA, foreground masking, "
    "weighted sampling. 2D models train for 100K steps, 3D for 300K. All on single L40S GPUs.\n\n"
    "What\u2019s different this time? First, proper validation from epoch 1 \u2014 we now compute "
    "per-class Dice for all 48 organelle classes during validation, with proper NaN masking "
    "for partially annotated crops. Everything is logged to TensorBoard.\n\n"
    "Second, the best model is now saved by validation Dice, not loss. This directly "
    "optimizes what the leaderboard measures.\n\n"
    "Third, we fixed a critical checkpoint resume bug: the best Dice score was being "
    "reset to -1 after every checkpoint restore, meaning that if a job got preempted and "
    "resumed, it would overwrite the best model with whatever came first after restart.\n\n"
    "Fourth, the 3D memory leak is fixed. CellMapDataLoader.refresh() recreates the "
    "DataLoader each epoch. With persistent workers enabled, old worker processes stay "
    "alive and accumulate TensorStore chunk cache, causing host RAM to balloon to 500+ GB "
    "within the first epoch. Setting persistent_workers=false lets workers die on refresh, "
    "releasing the cache. This is the fix we\u2019re testing now \u2014 3D jobs were submitted first "
    "so we can verify it early."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9: WHAT'S NEXT
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("What\u2019s Next")

# Timeline
tf = add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4))
set_para(tf, "Estimated Timeline", size=20, bold=True, color=BLUE)

timeline = [
    ("Now", "Phase 2 launched", "All 9 models + TensorBoard submitted to L40S",
     "\u2705 Running", GREEN),
    ("~1 hour", "3D OOM check", "Verify persistent_workers fix \u2014 3D jobs submitted first",
     "", GRAY),
    ("~6 hours", "First 2D Dice", "2D models hit epoch 10 \u2192 first per-class Dice scores",
     "", GRAY),
    ("~2 days", "2D complete", "All four 2D models finish 100 epochs (100K steps)",
     "", GRAY),
    ("~10\u201313 days", "3D complete", "All five 3D models finish 1000 epochs (300K steps)",
     "", GRAY),
    ("+ 2 days", "Threshold tuning", "Per-class sigmoid optimization on validation crops",
     "", GRAY),
    ("+ 1 day", "Per-class ensemble", "Select best architecture per organelle class",
     "", GRAY),
    ("+ 1 day", "Leaderboard", "Submit to CellMap challenge \u2192 objective score",
     "", GRAY),
]

for i, (when, title, desc, status, clr) in enumerate(timeline):
    y = Inches(1.35) + Inches(i * 0.62)
    bg = LTGREEN if i == 0 else (LTBLUE if i <= 4 else ZEBRA)

    # When label
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.5), y, Inches(1.5), Inches(0.50))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    tf2 = add_textbox(slide, Inches(0.55), y + Inches(0.05), Inches(1.4), Inches(0.40))
    set_para(tf2, when, size=12, bold=True, color=BLUE if i <= 4 else GRAY,
             alignment=PP_ALIGN.CENTER)

    # Title + description
    tf3 = add_textbox(slide, Inches(2.2), y + Inches(0.02), Inches(3.0), Inches(0.46))
    set_para(tf3, title, size=13, bold=True, color=DKGRAY)
    tf4 = add_textbox(slide, Inches(5.3), y + Inches(0.02), Inches(7.5), Inches(0.46))
    set_para(tf4, desc, size=11, color=GRAY)

# The vision box at bottom
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(6.55) - Inches(0.55), Inches(12.3), Inches(0.75))
rect.fill.solid()
rect.fill.fore_color.rgb = LTGREEN
rect.line.fill.background()

tf = add_textbox(slide, Inches(0.7), Inches(6.55) - Inches(0.48), Inches(11.9), Inches(0.6))
set_para(tf, "The Vision:  CellMap model  \u2192  per-class threshold tuning  \u2192  "
         "per-class ensemble  \u2192  apply to our own novel electron tomograms",
         size=14, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

set_notes(slide,
    "Here's the timeline going forward.\n\n"
    "We just submitted all 9 models to L40S GPUs. The 3D jobs were submitted first "
    "so we can verify the memory leak fix within the first hour. If host RAM stays "
    "under 100GB instead of ballooning to 500+, we know the fix works.\n\n"
    "First 2D Dice scores come at epoch 10 \u2014 about 6 hours in. 2D models should "
    "finish within 2 days. 3D models take 10 to 13 days for the full 300K steps.\n\n"
    "After training: per-class threshold tuning \u2014 the default 0.5 threshold is "
    "conservative, and precision typically far exceeds recall. Tuning per class "
    "could significantly boost Dice. Then per-class ensemble: different architectures "
    "may excel at different organelles, so the final model picks the best architecture "
    "for each of the 48 classes.\n\n"
    "The endgame: submit to the CellMap leaderboard for an objective score, then "
    "apply the trained model to our own novel electron tomograms \u2014 which was the "
    "whole point from the beginning."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10: THANK YOU
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(title_layout)

tf = add_textbox(slide, Inches(1), Inches(2.0), Inches(11.33), Inches(1.5))
set_para(tf, "Thank You!", size=48, bold=True, color=GOLD,
         alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

tf = add_textbox(slide, Inches(1), Inches(3.5), Inches(11.33), Inches(0.5))
set_para(tf, "Questions?", size=28, color=GRAY, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light")

tf = add_textbox(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(0.5))
set_para(tf, "George S. George  \u2022  gsgeorge@unc.edu", size=16, color=GRAY,
         alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, Inches(1), Inches(5.3), Inches(11.33), Inches(0.8))
set_para(tf, "Code: github.com/Shan-CU/CellMap-Segmentation", size=13, color=LTGRAY,
         alignment=PP_ALIGN.CENTER)

set_notes(slide, "Thank you! Happy to take any questions.")


# ── Save ────────────────────────────────────────────────────────────────
output_path = os.path.join(SCRIPT_DIR, "presentation_v3.pptx")
prs.save(output_path)
print(f"\n\u2705 Presentation v3 saved to: {output_path}")
print(f"   10 slides, Phase 2 focused")
print(f"   Speaker notes on all slides")
