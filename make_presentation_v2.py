#!/usr/bin/env python3
"""
make_presentation_v2.py — CellMap Lab Meeting Presentation (Revised Narrative)

9-slide narrative arc:
  1. Title
  2. The Big Picture / Why CellMap
  3. The Roadmap (8-step pipeline — native shapes)
  4. Choosing a Loss Function — quick test to pick a stable baseline
  5. The Pipeline + Ablation Design (merged)
  6. Ablation Design — 59 experiments, 5 sweeps
  7. 2D Results — Sweeps B & C
  8. 2D Results — Sweep D & Winners Summary
  9. NaN Debugging Story
 10. What's Next — status + remaining steps
 11. The Vision — CellMap → novel tomograms
 12. Thank You / Questions

Uses UNC lab template: GStephenson_Rotation3_Slides.pptx
Embeds LaTeX TikZ/pgfplots figures at 300 DPI.
"""

import io
import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# ── Paths ──────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE   = os.path.join(SCRIPT_DIR, "GStephenson_Rotation3_Slides.pptx")
FIG_DIR    = os.path.join(SCRIPT_DIR, "figures")

# ── Colour palette (matches template) ─────────────────────────────────
GOLD    = RGBColor(0xCF, 0xB8, 0x7C)
BLACK   = RGBColor(0x00, 0x00, 0x00)
DKGRAY  = RGBColor(0x33, 0x33, 0x33)
GRAY    = RGBColor(0x76, 0x76, 0x76)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x15, 0x60, 0x82)
ORANGE  = RGBColor(0xE9, 0x71, 0x32)
GREEN   = RGBColor(0x19, 0x6B, 0x24)
RED     = RGBColor(0xC0, 0x39, 0x2B)
LTGRAY  = RGBColor(0xD9, 0xD9, 0xD9)

# ── Helper functions ───────────────────────────────────────────────────

def fig_path(name):
    return os.path.join(FIG_DIR, name)


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_figure(slide, name, left, top, max_width, max_height):
    from PIL import Image
    path = fig_path(name)
    if not os.path.exists(path):
        print(f"  ⚠ Figure not found: {path}")
        return
    with Image.open(path) as img:
        w_px, h_px = img.size
    aspect = w_px / h_px
    w = max_width
    h = int(w / aspect)
    if h > max_height:
        h = max_height
        w = int(h * aspect)
    slide.shapes.add_picture(path, left, top, w, h)


def set_para(tf, text, size=18, bold=False, color=DKGRAY, alignment=PP_ALIGN.LEFT,
             space_after=Pt(4), space_before=Pt(0), font_name="Calibri"):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def add_para(tf, text, size=18, bold=False, color=DKGRAY, alignment=PP_ALIGN.LEFT,
             space_after=Pt(4), space_before=Pt(0), font_name="Calibri", level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    p.level = level
    return p


def add_bullet_list(tf, items, size=16, color=DKGRAY, bold=False,
                    space_after=Pt(6), level=0, bullet="\u2022"):
    for item in items:
        add_para(tf, f"{bullet} {item}", size=size, color=color, bold=bold,
                 space_after=space_after, level=level)


def add_table(slide, left, top, width, height, rows,
              header_bg=BLUE, header_fg=WHITE, font_size=11):
    n_rows = len(rows)
    n_cols = len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.font.name = "Calibri"
                if r_idx == 0:
                    para.font.bold = True
                    para.font.color.rgb = header_fg
                else:
                    para.font.color.rgb = DKGRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    return tbl


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ── Load template, extract logo, delete template slides ────────────────
_tpl = Presentation(TEMPLATE)

logo_blob = None
logo_ext  = "png"
for shape in _tpl.slides[1].shapes:
    if shape.name == "Picture 1":
        logo_blob = shape.image.blob
        logo_ext  = shape.image.content_type.split("/")[-1]
        break


def _delete_slide(prs, slide_index):
    rIds = []
    slide_part = prs.slides[slide_index].part
    for rId, rel in prs.part.rels.items():
        if rel.target_part is slide_part:
            rIds.append(rId)
    sldIdLst = prs.slides._sldIdLst
    ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    for sldId in list(sldIdLst):
        if sldId.get(ns + 'id') in rIds:
            sldIdLst.remove(sldId)
    for rId in rIds:
        prs.part.rels._rels.pop(rId, None)


for i in range(len(_tpl.slides) - 1, -1, -1):
    _delete_slide(_tpl, i)

prs = _tpl
title_layout = prs.slide_layouts[0]


def add_content_slide(title_text):
    slide = prs.slides.add_slide(title_layout)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0),
                                   Inches(13.33), Inches(6.50))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.fill.background()
    tf = add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.55))
    set_para(tf, title_text, size=28, bold=True, color=GOLD, font_name="Calibri Light")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), Inches(0.67),
                                   Inches(12.33), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    if logo_blob:
        buf = io.BytesIO(logo_blob)
        slide.shapes.add_picture(buf, Inches(11.45), Inches(6.57),
                                 Inches(1.88), Inches(0.86))
    return slide


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(title_layout)

tf = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10.93), Inches(2.0))
set_para(tf, "Multi-Label Organelle Segmentation", size=42, bold=True,
         color=GOLD, alignment=PP_ALIGN.LEFT, font_name="Calibri Light")
add_para(tf, "from FIB-SEM Electron Microscopy", size=42, bold=True,
         color=GOLD, alignment=PP_ALIGN.LEFT, font_name="Calibri Light")

tf = add_textbox(slide, Inches(1.2), Inches(3.8), Inches(10.93), Inches(0.6))
set_para(tf, "Toward Automated Annotation of Novel Electron Tomograms",
         size=22, color=DKGRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri")

tf = add_textbox(slide, Inches(1.2), Inches(5.0), Inches(10.93), Inches(0.5))
set_para(tf, "George S. George", size=24, bold=False, color=DKGRAY,
         alignment=PP_ALIGN.LEFT, font_name="Calibri")

tf = add_textbox(slide, Inches(1.2), Inches(5.5), Inches(10.93), Inches(0.5))
set_para(tf, "Department of Computer Science  \u2022  Interdisciplinary Quantitative Biology Program",
         size=14, color=GRAY, alignment=PP_ALIGN.LEFT)

tf = add_textbox(slide, Inches(1.2), Inches(5.85), Inches(10.93), Inches(0.3))
set_para(tf, "Lab Meeting  \u2022  February 2026", size=14, color=GRAY,
         alignment=PP_ALIGN.LEFT)

set_notes(slide,
    "Hi everyone, thanks for having me. "
    "Today I'll walk you through my rotation project \u2014 building a segmentation model "
    "for electron microscopy data. The long-term goal is a model that can take novel "
    "electron tomograms and automatically label every organelle.\n\n"
    "I'll cover why we chose CellMap as our training ground, the 8-step pipeline "
    "from baseline selection to novel tomograms, how we picked a stable loss function, "
    "the systematic ablation study \u2014 59 experiments across 5 sweeps \u2014 "
    "key 2D results so far, and where we're headed next."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE BIG PICTURE — WHY CELLMAP
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("The Big Picture: Why CellMap?")

# Left column
tf = add_textbox(slide, Inches(0.5), Inches(0.85), Inches(6.0), Inches(5.3))
set_para(tf, "The Long-Term Goal", size=20, bold=True, color=BLUE)
add_bullet_list(tf, [
    "Build a model that takes novel electron\n"
    "tomograms and produces fully labeled\n"
    "organelle segmentations \u2014 automatically",
    "No manual annotation \u2014 the model\n"
    "generalizes to unseen biology",
    "Everything we learn here transfers\n"
    "directly to our own tomograms",
], size=14, color=DKGRAY)

add_para(tf, "", size=10)
add_para(tf, "Why Is This Hard?", size=20, bold=True, color=BLUE, space_before=Pt(8))
add_bullet_list(tf, [
    "48 organelle classes, 22 biological samples",
    "Partial annotations \u2014 unlabeled \u2260 absent;\n"
    "standard losses penalize correct predictions",
    "6+ orders of magnitude class imbalance",
    "3D volumes need up to 512 GB host RAM",
], size=14, color=DKGRAY)

# Right: CellMap proving ground box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7.0), Inches(0.85), Inches(5.8), Inches(2.5))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE6)
rect.line.fill.background()

tf2 = add_textbox(slide, Inches(7.2), Inches(0.95), Inches(5.4), Inches(2.3))
set_para(tf2, "Why CellMap as the Training Ground?", size=16, bold=True, color=GOLD)
add_bullet_list(tf2, [
    "Largest public FIB-SEM segmentation benchmark",
    "Diverse biology: HeLa, mouse kidney/liver/\n"
    "heart, fly brain, jurkat, macrophage...",
    "Has a leaderboard \u2192 objective evaluation",
    "Same imaging modality as our own data",
], size=13, color=DKGRAY)

# Key stats boxes
data = [
    ("48", "organelle classes"),
    ("22", "biological samples"),
    ("289", "annotated 3D crops"),
    ("42.3B", "total voxels"),
]
for i, (num, label) in enumerate(data):
    y = Inches(3.7) + Inches(i * 0.75)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.2), y, Inches(5.4), Inches(0.65))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE8, 0xE0, 0xD0)
    rect.line.fill.background()
    tf3 = add_textbox(slide, Inches(7.4), y + Inches(0.05), Inches(1.5), Inches(0.55))
    set_para(tf3, num, size=26, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    tf4 = add_textbox(slide, Inches(8.8), y + Inches(0.12), Inches(3.6), Inches(0.4))
    set_para(tf4, label, size=14, color=GRAY)

set_notes(slide,
    "The long-term goal is to build a model that takes in novel electron tomograms \u2014 "
    "data our lab generates \u2014 and automatically produces fully labeled organelle "
    "segmentations. No manual annotation.\n\n"
    "Why CellMap? It's the largest publicly available FIB-SEM segmentation benchmark. "
    "48 organelle classes, 22 biological samples. It has a leaderboard so we can "
    "objectively measure how well we're doing. And it's the same imaging modality "
    "as our own tomograms.\n\n"
    "What makes this hard: annotations are partial \u2014 each crop only labels a subset "
    "of classes. Unlabeled does NOT mean absent. Class imbalance spans 6+ orders of "
    "magnitude. And 3D training needs 512 GB of host RAM just for data loading."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE ROADMAP (native python-pptx shapes)
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("The Roadmap")

steps = [
    ("1", "Baseline Selection",       "Quick test to identify stable loss function",         "done"),
    ("2", "Infrastructure Rebuild",   "Native zarr pipeline, all 48 classes",                  "done"),
    ("3", "Systematic Ablation",      "59 experiments, 5 sweeps (2D + 3D)",                    "here"),
    ("4", "Compose Optimal Config",   "Best of each sweep \u2192 single config",              ""),
    ("5", "Full Architecture Training","8 architectures \u00d7 100 epochs",                    ""),
    ("6", "Threshold Tuning",         "Per-class IoU optimization on validation",               ""),
    ("7", "Per-Class Ensemble",       "Best architecture per organelle",                        ""),
    ("8", "Apply to Novel Tomograms", "The Goal",                                               "goal"),
]

y_start = Inches(0.90)
row_h   = Inches(0.62)
gap     = Inches(0.06)

for i, (num, title, desc, status) in enumerate(steps):
    y = y_start + (row_h + gap) * i

    # Background colors
    if status == "done":
        bg = RGBColor(0xE8, 0xE0, 0xD0)
        num_bg = GOLD
    elif status == "here":
        bg = RGBColor(0xD0, 0xDE, 0xE8)
        num_bg = BLUE
    elif status == "goal":
        bg = RGBColor(0xD4, 0xE8, 0xD4)
        num_bg = GOLD
    else:
        bg = RGBColor(0xF2, 0xF2, 0xF2)
        num_bg = LTGRAY

    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(0.65), y + Inches(0.06),
                                   Inches(0.48), Inches(0.48))
    circ.fill.solid()
    circ.fill.fore_color.rgb = num_bg
    circ.line.fill.background()
    tf_n = add_textbox(slide, Inches(0.65), y + Inches(0.09),
                        Inches(0.48), Inches(0.44))
    set_para(tf_n, num, size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Step box
    box_w = Inches(3.8)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.30), y, box_w, row_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    if status == "here":
        rect.line.color.rgb = BLUE
        rect.line.width = Pt(2)

    # Step title
    tf_t = add_textbox(slide, Inches(1.50), y + Inches(0.08), Inches(3.4), Inches(0.44))
    title_color = BLUE if status == "here" else DKGRAY
    set_para(tf_t, title, size=15, bold=True, color=title_color)

    # Description
    tf_d = add_textbox(slide, Inches(5.30), y + Inches(0.08), Inches(5.0), Inches(0.44))
    set_para(tf_d, desc, size=13, color=GRAY)

    # Status label
    if status == "done":
        tf_s = add_textbox(slide, Inches(10.5), y + Inches(0.08), Inches(2.3), Inches(0.44))
        set_para(tf_s, "\u2705 Done", size=14, bold=True, color=GREEN,
                 alignment=PP_ALIGN.RIGHT)
    elif status == "here":
        tf_s = add_textbox(slide, Inches(10.5), y + Inches(0.08), Inches(2.3), Inches(0.44))
        set_para(tf_s, "\u2190 In Progress", size=14, bold=True, color=BLUE,
                 alignment=PP_ALIGN.RIGHT)

    # Arrow to next step (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                        Inches(3.0), y + row_h,
                                        Inches(0.30), gap)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

set_notes(slide,
    "Here's the full pipeline from start to finish. We're at step 3.\n\n"
    "Step 1: ran a short test to identify a stable loss function for ablations. "
    "Step 2: rebuilt the entire pipeline with native zarr and all 48 classes. "
    "Step 3 is where we are now \u2014 59 experiments across 5 sweeps.\n\n"
    "Steps 4\u20138 are ahead: compose optimal config, full architecture training, "
    "per-class threshold tuning, per-class ensemble, and apply to novel tomograms."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: CHOOSING A LOSS FUNCTION
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("Step 1: Choosing a Baseline Loss Function")

# Central message — clean and spacious
tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.0))
set_para(tf, "Before running 59 ablation experiments, we needed a stable loss function "
         "that would give real, comparable results across sweeps.",
         size=20, color=DKGRAY)

# The approach box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.0))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE6)
rect.line.fill.background()

tf = add_textbox(slide, Inches(1.8), Inches(2.65), Inches(9.7), Inches(1.8))
set_para(tf, "What we did", size=18, bold=True, color=GOLD)
add_bullet_list(tf, [
    "Ran a short test across several common loss functions",
    "Identified BalancedSoftmaxTversky (BST) as a stable baseline \u2014 "
    "handles partial annotations and class imbalance without collapsing",
    "Used BST as the constant loss in all ablation sweeps, "
    "isolating the effect of each variable being tested",
], size=15, color=DKGRAY, space_after=Pt(10))

# Outcome callout
rect2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(3.0), Inches(4.9), Inches(7.3), Inches(0.75))
rect2.fill.solid()
rect2.fill.fore_color.rgb = RGBColor(0xD0, 0xDE, 0xE8)
rect2.line.fill.background()

tf = add_textbox(slide, Inches(3.2), Inches(4.95), Inches(6.9), Inches(0.65))
set_para(tf, "\u2192  BST loss = constant across all sweeps \u2192 minimal confounding variables",
         size=16, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

set_notes(slide,
    "Before spending thousands of GPU hours on ablations, we ran a short test "
    "to find a loss function that was stable and gave real results.\n\n"
    "BalancedSoftmaxTversky handles partial annotations and 6 orders of magnitude "
    "of class imbalance without collapsing. It became the constant loss "
    "across all our ablation sweeps, so each sweep isolates just one variable."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5: THE PIPELINE + ABLATION DESIGN (merged)
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("The Pipeline \u2014 59 Ablation Experiments")

# Pipeline visual (top)
add_figure(slide, "pipeline_visual.png",
           left=Inches(0.15), top=Inches(0.80),
           max_width=Inches(13.0), max_height=Inches(2.8))

# Sweep table (bottom)
add_figure(slide, "sweep_table.png",
           left=Inches(0.3), top=Inches(3.75),
           max_width=Inches(12.7), max_height=Inches(2.7))

set_notes(slide,
    "Here's the full pipeline. Data flows left to right: "
    "22 zarr volumes, 289 crops across 48 classes, "
    "weighted sampler, FlexUNet with ResNet-34 encoder.\n\n"
    "The loss pipeline: NaN mask for partial annotations, "
    "BalancedSoftmaxTversky loss, and logit adjustment.\n\n"
    "5 ablation sweeps, each varying one component. "
    "Sweep A: 8 loss functions. B: Tversky alpha-beta (6). "
    "C: class weight tau (5). D: masking strategy (7). E: training tricks (3+).\n\n"
    "59 experiments total: 29 in 2D, 30 in 3D. "
    "2D takes ~3h on L40S, 3D takes ~15h on H100 with 512 GB RAM."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6: 2D RESULTS — SWEEPS B & C
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("2D Results: Tversky \u03b1/\u03b2 and Class Weighting \u03c4")

tf = add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.25))
set_para(tf, "Val losses comparable within each sweep (same loss family). 25/29 complete; 4 relaunched after bug fixes.",
         size=11, color=GRAY)

# Sweep B chart (left)
add_figure(slide, "sweep_b_chart.png",
           left=Inches(0.2), top=Inches(1.15),
           max_width=Inches(6.3), max_height=Inches(3.6))

# Sweep C chart (right)
add_figure(slide, "sweep_c_chart.png",
           left=Inches(6.8), top=Inches(1.15),
           max_width=Inches(6.0), max_height=Inches(3.6))

# Winner B
tf = add_textbox(slide, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.3))
set_para(tf, "Winner: \u03b1=0.6, \u03b2=0.4  (val_loss = 0.561)", size=14, bold=True, color=BLUE)
add_para(tf, "Mild precision bias outperforms\n"
         "balanced Dice (\u03b1=\u03b2=0.5).",
         size=12, color=GRAY)

# Winner C
tf = add_textbox(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.3))
set_para(tf, "Winner: \u03c4=2.0  (val_loss = 0.409)", size=14, bold=True, color=BLUE)
add_para(tf, "Stronger logit adjustment needed\n"
         "for 48 classes with extreme imbalance.",
         size=12, color=GRAY)

set_notes(slide,
    "Sweep B: Tversky alpha-beta tradeoff. Alpha 0.6, beta 0.4 wins at 0.561. "
    "Mild precision bias \u2014 penalizing false positives more \u2014 works best.\n\n"
    "Sweep C: logit adjustment tau. Tau 2.0 wins. "
    "With 48 classes and extreme imbalance, a stronger logit adjustment "
    "is needed to prevent the majority classes from dominating."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: 2D RESULTS — SWEEP D & WINNERS SUMMARY
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("2D Results: Masking Strategy & Winners Summary")

# Sweep D chart (left, more room)
add_figure(slide, "sweep_d_chart.png",
           left=Inches(0.1), top=Inches(0.85),
           max_width=Inches(7.8), max_height=Inches(5.5))

# Winners summary panel (right)
tf = add_textbox(slide, Inches(8.2), Inches(0.85), Inches(4.8), Inches(5.5))
set_para(tf, "All Sweep Winners", size=20, bold=True, color=BLUE)

add_para(tf, "", size=8)
add_para(tf, "B: \u03b1=0.6, \u03b2=0.4  (0.561)", size=14, bold=True, color=GREEN)
add_para(tf, "Mild precision bias outperforms balanced Dice", size=12, color=GRAY)

add_para(tf, "", size=8)
add_para(tf, "C: \u03c4=2.0  (0.409)", size=14, bold=True, color=GREEN)
add_para(tf, "Stronger logit adjustment needed for 48 classes", size=12, color=GRAY)

add_para(tf, "", size=8)
add_para(tf, "D: masksup + FG, no bbox  (0.384)", size=14, bold=True, color=GREEN)
add_para(tf, "Surprise: tight bbox alone HURTS (0.719)\n"
         "FG mask + mask-supervised recon wins", size=12, color=GRAY)

add_para(tf, "", size=8)
add_para(tf, "A: BST  (pending 2 relaunches)", size=14, bold=True, color=ORANGE)
add_para(tf, "focal_tversky & unified_focal still running", size=12, color=GRAY)

add_para(tf, "", size=8)
add_para(tf, "E: EMA leading  (pending 2 relaunches)", size=14, bold=True, color=ORANGE)
add_para(tf, "no_weighted_sampler & focal_tversky_mild still running", size=12, color=GRAY)

add_para(tf, "", size=14)
# Composed config box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(8.2), Inches(5.0), Inches(4.8), Inches(1.2))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xD0, 0xDE, 0xE8)
rect.line.fill.background()

tf2 = add_textbox(slide, Inches(8.4), Inches(5.05), Inches(4.4), Inches(1.1))
set_para(tf2, "Preliminary Phase 2 Config:", size=12, bold=True, color=BLUE)
add_para(tf2, "BST(\u03c4=2.0, \u03b1=0.6, \u03b2=0.4)\n"
         "+ masksup=0.3, no bbox, FG mask\n"
         "+ EMA(0.999)  [pending final A & E]",
         size=11, color=DKGRAY, font_name="Consolas")

set_notes(slide,
    "The masking sweep is the most interesting. "
    "No masking gives 0.678. FG masking alone: 0.487 \u2014 28% improvement for free. "
    "But tight bbox alone is WORSE than nothing at 0.719.\n\n"
    "Winner: mask-supervised reconstruction plus FG masking without bbox at 0.384.\n\n"
    "Putting it all together: BST with tau=2.0, alpha=0.6, beta=0.4, "
    "masksup=0.3, no bbox, foreground mask, and EMA. "
    "That becomes the Phase 2 baseline config."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: WHAT'S NEXT
# ════════════════════════════════════════════════════════════════════════
slide = add_content_slide("What\u2019s Next: Steps 4\u20138")

# Status table
status_rows = [
    ["Component", "Status", "Details"],
    ["2D Ablation (25/29)", "\u2705 Complete", "All converged at 50 epochs"],
    ["2D Relaunches (4)", "\U0001f7e1 Running", "Bug fix applied \u2192 relaunched on L40S"],
    ["3D Ablation (11/30)", "\U0001f7e1 Running", "Sycamore H100 cluster"],
    ["3D Pending (22/30)", "\u23f3 Queued", "~15h per job \u2014 done ~Wed Feb 26"],
]
add_table(slide, Inches(0.5), Inches(0.85), Inches(6.0), Inches(2.0),
          status_rows, font_size=11, header_bg=BLUE)

# Hardware + endgame box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7.0), Inches(0.85), Inches(5.8), Inches(2.0))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE6)
rect.line.fill.background()

tf = add_textbox(slide, Inches(7.2), Inches(0.95), Inches(5.4), Inches(1.8))
set_para(tf, "Compute & Timeline", size=14, bold=True, color=GOLD)
add_para(tf, "", size=4)
add_para(tf, "2D relaunches: done within hours (L40S)", size=11, color=DKGRAY)
add_para(tf, "3D ablation: ~Wed Feb 26 (H100 cluster)", size=11, color=DKGRAY)
add_para(tf, "Phase 2 architecture training: ~2 weeks", size=11, color=DKGRAY)
add_para(tf, "~2,200 GPU-hours spent so far", size=11, color=GRAY)

# Pipeline ahead
tf = add_textbox(slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.4))
set_para(tf, "Remaining Pipeline", size=20, bold=True, color=BLUE)

next_steps = [
    ("4. Compose Config", "Take winner from each sweep \u2192 single optimal configuration",
     "~1 day after ablation"),
    ("5. Full Arch Training", "4 archs \u00d7 2D + 4 archs \u00d7 3D, 100 epochs, 1000 iter/epoch",
     "~2 weeks on cluster"),
    ("6. Threshold Tuning", "Per-class threshold optimization on validation (default 0.5 is suboptimal)",
     "~2 days"),
    ("7. Per-Class Ensemble", "Different architectures excel at different organelles \u2192 composite",
     "~1 day"),
    ("8. Final Model", "Submit to CellMap leaderboard \u2192 apply to novel tomograms",
     "End goal"),
]

for i, (title, desc, timeline) in enumerate(next_steps):
    y = Inches(3.6) + Inches(i * 0.58)
    bg_color = RGBColor(0xD4, 0xE8, 0xD4) if i == 4 else RGBColor(0xF2, 0xF2, 0xF2)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.5), y, Inches(8.5), Inches(0.48))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.fill.background()

    tf2 = add_textbox(slide, Inches(0.7), y + Inches(0.02), Inches(2.3), Inches(0.42))
    set_para(tf2, title, size=12, bold=True, color=DKGRAY)
    tf3 = add_textbox(slide, Inches(3.0), y + Inches(0.02), Inches(5.7), Inches(0.42))
    set_para(tf3, desc, size=10, color=GRAY)
    tf4 = add_textbox(slide, Inches(9.2), y + Inches(0.02), Inches(3.5), Inches(0.42))
    set_para(tf4, timeline, size=10, bold=True,
             color=GREEN if i == 4 else GRAY, alignment=PP_ALIGN.RIGHT)

set_notes(slide,
    "25 of 29 2D experiments complete, 4 relaunched after a bug fix. "
    "3D ablation has 11 of 30 done, rest queued, should finish by Wednesday.\n\n"
    "Then: compose winning config, full architecture training across 8 models, "
    "per-class threshold tuning, per-class ensemble, submit to leaderboard, "
    "and apply to our own novel tomograms."
)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9: THANK YOU
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(title_layout)

tf = add_textbox(slide, Inches(1), Inches(2.0), Inches(11.33), Inches(1.5))
set_para(tf, "Thank You!", size=48, bold=True, color=GOLD,
         alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

tf = add_textbox(slide, Inches(1), Inches(3.5), Inches(11.33), Inches(0.5))
set_para(tf, "Questions?", size=28, color=GRAY, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light")

tf = add_textbox(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(0.5))
set_para(tf, "George S. George  \u2022  gsgeorge@unc.edu", size=16, color=GRAY,
         alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, Inches(1), Inches(5.3), Inches(11.33), Inches(0.8))
set_para(tf, "Code: github.com/Shan-CU/CellMap-Segmentation", size=13, color=LTGRAY,
         alignment=PP_ALIGN.CENTER)

set_notes(slide,
    "Thank you! Happy to take any questions."
)


# ── Save ────────────────────────────────────────────────────────────────
output_path = os.path.join(SCRIPT_DIR, "presentation_v2.pptx")
prs.save(output_path)
print(f"\n\u2705 Presentation v2 saved to: {output_path}")
print(f"   9 slides using template theme/layout")
print(f"   Speaker notes on all slides")
